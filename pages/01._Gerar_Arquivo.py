import streamlit as st
import mysql.connector
from docx import Document
from docx.shared import Pt
from docx.enum.text import WD_ALIGN_PARAGRAPH
import io
import openpyxl
import os
from pptx import Presentation
import re



tab1,tab2, tab3,tab4 = st.tabs(['Teoria', 'Questões', 'Slide Questões', 'Converter JPG'])

with tab1:
    # Funções auxiliares e a função de conversão
    def aplicar_substituicoes(texto, substituicoes):
        for palavra_antiga, palavra_nova in substituicoes:
            texto = texto.replace(palavra_antiga, palavra_nova)
        return texto

    def dividir_texto(texto, limite_caracteres, delimitador):
        partes = []
        parte_atual = ""

        for palavra in texto.split():
            if len(parte_atual + ' ' + palavra) > limite_caracteres and parte_atual.endswith(delimitador):
                partes.append(parte_atual)
                parte_atual = palavra
            else:
                if parte_atual:
                    parte_atual += ' '
                parte_atual += palavra

        if parte_atual:
            partes.append(parte_atual)

        return partes

    def pptx_to_word_with_slide_markers(pptx_memory, word_memory, substituicoes, limite_caracteres, delimitador):
        prs = Presentation(pptx_memory)
        doc = Document()

        for i, slide in enumerate(prs.slides):
            doc.add_paragraph(f"SLIDE: {i+1}")

            if i > 0:
                doc.add_paragraph().add_run().add_break()

            title_text = ""
            if slide.shapes.title:
                title_text = slide.shapes.title.text
                if title_text:
                    title_text = aplicar_substituicoes(title_text, substituicoes)
                    doc.add_paragraph("TITLE: " + title_text, style='Heading 1')

            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text and shape.text != title_text:
                    shape_text = aplicar_substituicoes(shape.text, substituicoes)
                    partes_texto = dividir_texto(shape_text, limite_caracteres, delimitador)
                    for parte in partes_texto:
                        doc.add_paragraph(parte)

        doc.save(word_memory)

    # Substituições de palavras e funções de apoio
    substituicoes = [
        ("Hely", "Elí"),
        ("Di Pietro", "Di piêtro"),
        ("CF", "Constituição Federal"),
        ("nº", "número"),
        ("art.", "artigo"),
        ("Art.", "artigo"),
        ("obs.", "observação"),
        ("Obs.:", "observação"),
        ("J.J.", "José Joaquim"),
        ("j.j.", "José Joaquim"),
        ("habeas", "habias"),
        ("Habeas", "habias"),
        ("corpus", "corpos"),
        ("§", "parágrafo"),
        ("LOA", "Lei Orçamentária Anual"),
    ]
    limite_caracteres = 100  # Limite de caracteres para divisão do texto
    delimitador = '.'        # Delimitador para divisão do texto

    def streamlit_app():
        st.title("Conversor PPTX para DOCX com Marcadores de Slide")

        pptx_file = st.file_uploader("Escolha o arquivo PPTX", type="pptx")
        if pptx_file is not None:
            pptx_memory = io.BytesIO(pptx_file.getvalue())
            word_memory = io.BytesIO()

            if st.button('Converter PPTX para DOCX'):
                with st.spinner('Convertendo...'):
                    # Iniciar a barra de progresso
                    progress_bar = st.progress(0)
                    for i in range(1, 101):
                        # Atualizando a barra de progresso
                        progress_bar.progress(i)
                    pptx_to_word_with_slide_markers(pptx_memory, word_memory, substituicoes, limite_caracteres, delimitador)
                    
                    # Resetar a barra de progresso após a conclusão
                    progress_bar.empty()

                # Exibir mensagem de sucesso
                st.success('Conversão concluída com sucesso!')

                word_memory.seek(0)
                st.download_button(
                    label="Baixar arquivo DOCX",
                    data=word_memory,
                    file_name='documento_convertido.docx',
                    mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document"
                )

    if __name__ == "__main__":
        streamlit_app()

with tab2:
    # Função para aplicar substituições de palavras
    def aplicar_substituicoes(texto, substituicoes):
        for palavra_antiga, palavra_nova in substituicoes:
            texto = texto.replace(palavra_antiga, palavra_nova)
        return texto

    # Função para extrair a parte do gabarito comentado
    def extrair_gabarito_comentado(comentario):
        indice = comentario.find("**Gabarito Comentado:**")
        if indice != -1:
            return comentario[indice + len("**Gabarito Comentado:**"):].strip()
        return ""


    def limpar_texto(texto):
        # Certifique-se de que o texto é uma string
        if not isinstance(texto, str):
            texto = str(texto)

        # Remove caracteres incompatíveis com XML
        texto = re.sub(r'[\x00-\x1F\x7F-\x9F]', '', texto)

        return texto

    # Função para adicionar conteúdo ao documento
    def adicionar_conteudo_ao_documento(doc, slide_numero, numero_da_questao, titulo, conteudo, substituicoes, incluir_titulo=True):
    # ... Restante da função ...
        p = doc.add_paragraph()
        run = p.add_run(f"SLIDE: {slide_numero}\n")
        run.bold = True
        run.font.size = Pt(16)
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER

        conteudo_formatado = limpar_texto(conteudo)
        if incluir_titulo:
            titulo_formatado = aplicar_substituicoes(f"{titulo} {numero_da_questao}", substituicoes)
            doc.add_paragraph(titulo_formatado, style='Heading 1')

        conteudo_formatado = aplicar_substituicoes(conteudo, substituicoes)
        doc.add_paragraph(conteudo_formatado)

    def gerar_documento_com_questoes(materia, assunto, topico, config, substituicoes):
        doc = Document()
        slide_numero = 1
        numero_da_questao = 1
        questoes_obtidas = []
        
        if assunto:
            titulo_assunto = f"Questões Comentadas sobre {assunto}"
            p = doc.add_paragraph()
            run = p.add_run(f"SLIDE: {slide_numero}\n" + titulo_assunto)
            
            run.bold = True
            run.font.size = Pt(16)
            p.alignment = WD_ALIGN_PARAGRAPH.CENTER
            slide_numero += 1  # Incrementa o número do slide

        try:
            cnx = mysql.connector.connect(**config)
            cursor = cnx.cursor()

            # Construa sua consulta SQL aqui. A consulta abaixo é um exemplo e pode precisar ser ajustada.
            base_query = """
                SELECT Q.QuestaoID, Q.Questao, Q.Comentario
                FROM Questoes Q
                JOIN Materias M ON Q.MateriaID = M.MateriaID
                JOIN Assuntos A ON Q.AssuntoID = A.AssuntoID
                JOIN Topicos T ON Q.TopicoID = T.TopicoID
                WHERE (%s = '' OR M.Materia = %s)
                AND (%s = '' OR A.Assunto = %s)
                AND (%s = '' OR T.Topico = %s)
            """

            cursor.execute(base_query, (materia, materia, assunto, assunto, topico, topico))
            questoes = cursor.fetchall()

            for questaoID, questao, comentario in questoes:
                adicionar_conteudo_ao_documento(doc, slide_numero, numero_da_questao, "Questão", questao, substituicoes)
                slide_numero += 1

                gabarito_comentado = extrair_gabarito_comentado(comentario)
                if gabarito_comentado:
                    adicionar_conteudo_ao_documento(doc, slide_numero, numero_da_questao, "", gabarito_comentado, substituicoes, incluir_titulo=False)
                    slide_numero += 1

                questoes_obtidas.append((questaoID, questao, comentario))
                numero_da_questao += 1
                
        except mysql.connector.Error as err:
            print(f"Erro ao conectar ao MySQL: {err}")
        finally:
            if cnx.is_connected():
                cursor.close()
                cnx.close()

        return doc
    # Função para salvar informações em uma planilha Excel

    # Função para obter matérias
    def obter_materias(config):
        try:
            cnx = mysql.connector.connect(**config)
            cursor = cnx.cursor()

            cursor.execute("SELECT Materia FROM Materias;")
            materias = [str(row[0]) for row in cursor.fetchall()]
            return materias

        except Exception as e:
            st.error(f"Erro ao obter matérias: {e}")
            return []

        finally:
            if cnx.is_connected():
                cursor.close()
                cnx.close()

    # Função para obter assuntos com base na matéria
    def obter_assuntos(config, materia):
        try:
            cnx = mysql.connector.connect(**config)
            cursor = cnx.cursor()

            cursor.execute("""
                SELECT DISTINCT Assunto 
                FROM Assuntos
                JOIN Questoes ON Assuntos.AssuntoID = Questoes.AssuntoID
                JOIN Materias ON Questoes.MateriaID = Materias.MateriaID
                WHERE Materias.Materia = %s
                GROUP BY Assuntos.Assunto;
            """, (materia,))
            assuntos = [str(row[0]) for row in cursor.fetchall()]
            return assuntos

        except Exception as e:
            st.error(f"Erro ao obter assuntos: {e}")
            return []

        finally:
            if cnx.is_connected():
                cursor.close()
                cnx.close()

    # Função para obter tópicos com base na matéria e assunto
    def obter_topicos(config, materia, assunto):
        try:
            cnx = mysql.connector.connect(**config)
            cursor = cnx.cursor()

            cursor.execute("""
                SELECT DISTINCT Topico 
                FROM Topicos
                JOIN Questoes ON Topicos.TopicoID = Questoes.TopicoID
                JOIN Assuntos ON Questoes.AssuntoID = Assuntos.AssuntoID
                JOIN Materias ON Questoes.MateriaID = Materias.MateriaID
                WHERE Materias.Materia = %s AND Assuntos.Assunto = %s
                GROUP BY Topicos.Topico;
            """, (materia, assunto))
            topicos = [str(row[0]) for row in cursor.fetchall()]
            return topicos

        except Exception as e:
            st.error(f"Erro ao obter tópicos: {e}")
            return []

        finally:
            if cnx.is_connected():
                cursor.close()
                cnx.close()

    def streamlit_app():
        st.title("Gerador de Documentos de Questões")

        # Configurações de conexão ao banco de dados
        config = {
            'user': 'admin',
            'password': 'Eduardo13*',
            'host': 'institutoscheffelt.ckrs9teerzcf.sa-east-1.rds.amazonaws.com',
            'database': 'questoes',
            'raise_on_warnings': True
        }

        # Substituições de palavras desejadas
        substituicoes = [
            ("Hely", "Elí"),
            ("Di Pietro", "Di piêtro"),
            ("CF","Constituição Federal"),
            ("nº","número"),
            ("n.º","número"),
            ("art.","artigo"),
            ("J.J.","José Joaquim"),
            ("habeas","habias"),
            ("corpus","corpos"),
            ("§", "parágrafo"),
            ("LOA", "Lei Orçamentária Anual"),
            ("RREO", "RÊO")
            # Adicione outras substituições conforme necessário
        ]
    # Seleção de Matéria
        materias = obter_materias(config)
        materia_selecionada = st.selectbox("Escolha a Matéria:", [""] + materias)

        # Seleção de Assunto
        if materia_selecionada:
            assuntos = obter_assuntos(config,materia_selecionada)
            assunto_selecionado = st.selectbox("Escolha o Assunto:", [""] + assuntos)
        else:
            assunto_selecionado = ""

        # Seleção de Tópico
        if assunto_selecionado and materia_selecionada:
            topicos = obter_topicos(config,materia_selecionada, assunto_selecionado)
            topico_selecionado = st.selectbox("Escolha o Tópico:", [""] + topicos)
        else:
            topico_selecionado = ""


        # Botão para gerar documento
        if st.button("Gerar Documento"):
            with st.spinner("Gerando..."):
                documento = gerar_documento_com_questoes(materia_selecionada, assunto_selecionado, topico_selecionado, config, substituicoes)

                buffer = io.BytesIO()
                documento.save(buffer)
                buffer.seek(0)


                st.download_button(label="Baixar Documento",
                                data=buffer,
                                file_name="questoes.docx",
                                mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document")

    if __name__ == "__main__":
        streamlit_app()

with tab3:
    import streamlit as st
    import io
    from pptx import Presentation
    from pptx.util import Pt
    from pptx.enum.text import PP_ALIGN

    def center_last_textbox_vertically(slide, prs_height, font_size, alignment):
        last_textbox = None

        # Identificar a última caixa de texto
        for shape in reversed(slide.shapes):
            if shape.has_text_frame:
                last_textbox = shape
                break

        if last_textbox:
            # Centralizar horizontalmente
            text_frame = last_textbox.text_frame
            text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER if alignment == 'center' else PP_ALIGN.LEFT

            # Ajustar tamanho da fonte da última caixa de texto
            for paragraph in text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(font_size)

            # Centralizar verticalmente
            text_height = last_textbox.height
            vertical_position = int((prs_height - text_height) / 2)
            last_textbox.top = vertical_position

    def modify_and_download_slides(presentation_path, font_size, alignment, keyword="Comentário:"):
        # Carregar a apresentação
        prs = Presentation(presentation_path)

        # Iterar por todos os slides
        for i, slide in enumerate(prs.slides):
            shapes_to_delete = []

            # Identificar caixas de texto para exclusão
            for j, shape in enumerate(slide.shapes):
                if shape.has_text_frame and keyword.lower() in shape.text_frame.text.lower():
                    # Marcar esta e a próxima caixa de texto para exclusão
                    shapes_to_delete.extend([shape, slide.shapes[j + 1]] if j + 1 < len(slide.shapes) else [shape])

            # Excluir as caixas de texto marcadas
            for shape in shapes_to_delete:
                sp = shape._element
                sp.getparent().remove(sp)

            # Verificar se caixas de texto foram excluídas neste slide
            if shapes_to_delete:
                # Centralizar a última caixa de texto restante apenas no slide atual
                center_last_textbox_vertically(slide, prs.slide_height, font_size, alignment)


        # Salvar a apresentação modificada
        #prs.save('sem_comentario.pptx')

        pptx_bytes = io.BytesIO()
        prs.save(pptx_bytes)
        pptx_bytes.seek(0)  # Posicionar o cursor no início do fluxo de bytes

        return pptx_bytes

    # Configurações da aplicação Streamlit
    st.title("Modificação de Apresentação PowerPoint")
    presentation_path = st.file_uploader("Selecione sua apresentação PowerPoint (.pptx)", type="pptx")
    font_size = st.number_input("Novo tamanho da fonte:", min_value=1, step=1, value=24)
    alignment = st.selectbox("Alinhamento:", options=['center', 'left'], index=0)

    if presentation_path is not None:
        pptx_bytes = None
        if st.button('Modificar e Baixar Slides'):
            pptx_bytes = modify_and_download_slides(presentation_path, font_size, alignment)

            # Baixar a apresentação modificada
            st.download_button(
                label="Baixar arquivo PPTX",
                data=pptx_bytes,
                file_name='sem_comentario.pptx',
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
            )

with tab4:
    import streamlit as st
    from pptx import Presentation
    from PIL import Image, ImageDraw

    def pptx_to_images(input_pptx, output_format='JPEG'):
        # Abre o arquivo PowerPoint
        prs = Presentation(input_pptx)

        # Itera sobre cada slide no PowerPoint
        for i, slide in enumerate(prs.slides):
            # Cria uma nova imagem com as dimensões do slide
            img = Image.new("RGB", (prs.slide_width, prs.slide_height), "white")
            img_draw = ImageDraw.Draw(img)

            # Renderiza o slide na imagem
            slide.shapes._spTree.render(img_draw)

            # Salva a imagem no formato desejado (JPEG ou PNG)
            img.save(f"slide_{i}.{output_format.lower()}", output_format)

    # Interface do Streamlit
    st.title("Conversor PPTX para Imagens")

    # Upload do arquivo PowerPoint
    uploaded_file = st.file_uploader("Faça o upload do arquivo PowerPoint (.pptx)", type="pptx")

    if uploaded_file:
        # Exibir informações sobre o arquivo carregado
        st.write("Arquivo carregado:", uploaded_file.name)
        
        # Botão para converter o arquivo
        if st.button("Converter para Imagens"):
            # Converter o arquivo PPTX em imagens
            pptx_to_images(uploaded_file, output_format='JPEG')
            st.success("Conversão concluída. Verifique as imagens geradas.")


